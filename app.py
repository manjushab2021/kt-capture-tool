"""
Knowledge Transfer Capture Tool
=================================
Built with: Python, Flask, Groq API, PyPDF2, python-docx
Author: Manjusha Bhakta

Run:
    pip install flask groq httpx pypdf2 python-docx
    set GROQ_API_KEY=gsk_your-key
    python app.py
    Open: http://localhost:5000
"""

import os
import io
import json
import httpx
from datetime import datetime
from flask import Flask, render_template, request, jsonify, send_file
from groq import Groq
import PyPDF2
from docx import Document as WordDoc
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

app = Flask(__name__)
client = Groq(
    api_key=os.environ.get("GROQ_API_KEY", "gsk_your-key-here"),
    http_client=httpx.Client()
)

# In-memory store
store = {
    "documents": [],       # list of {filename, text}
    "questions": [],       # list of {category, question}
    "answers": {},         # {question_index: answer}
    "person_name": "",
    "role": "",
    "project": ""
}


def extract_text(file):
    try:
        filename = file.filename.lower()
        content = file.read()

        if filename.endswith(".pdf"):
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text if text else "Could not extract text from PDF"

        elif filename.endswith(".txt"):
            try:
                return content.decode("utf-8")
            except:
                return content.decode("latin-1")

        elif filename.endswith(".docx"):
            from docx import Document as WordDoc
            doc = WordDoc(io.BytesIO(content))
            text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            return text if text else "Could not extract text from Word document"

        elif filename.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text if text else "Could not extract text from PowerPoint"

        else:
            return None

    except Exception as e:
        return f"Error reading file: {str(e)}"

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/upload", methods=["POST"])
def upload():
    files = request.files.getlist("files")
    uploaded = []

    for file in files:
        text = extract_text(file)
        if text:
            store["documents"].append({
                "filename": file.filename,
                "text": text[:6000]  # limit per doc
            })
            uploaded.append(file.filename)

    return jsonify({
        "success": True,
        "uploaded": uploaded,
        "total_docs": len(store["documents"])
    })


@app.route("/generate-questions", methods=["POST"])
def generate_questions():
    data = request.json
    store["person_name"] = data.get("person_name", "")
    store["role"] = data.get("role", "")
    store["project"] = data.get("project", "")

    if not store["documents"]:
        return jsonify({"error": "No documents uploaded yet"}), 400

    # Combine all document text
    combined = ""
    for doc in store["documents"]:
        combined += f"\n\n=== {doc['filename']} ===\n{doc['text']}"

    prompt = f"""You are an expert in knowledge transfer for data and analytics teams.

A {store['role'] or 'Data/Business Analyst'} named {store['person_name'] or 'the analyst'} is leaving the organisation.
Project/context: {store['project'] or 'Data and analytics program'}

Based on the documents below, generate exactly 15 tailored knowledge transfer questions.
The questions should be specific to what is in these documents — not generic.

Cover these categories (2-3 questions each):
1. Project Context & Objectives
2. Current State & Target State
3. Key Decisions & Rationale
4. Data Sources & Data Quality
5. Stakeholders & Relationships
6. Risks, Issues & Known Gaps
7. Handover Advice

Return ONLY a JSON array in this exact format, no other text:
[
  {{"category": "Project Context & Objectives", "question": "specific question here"}},
  {{"category": "Current State & Target State", "question": "specific question here"}}
]

Documents:
{combined[:8000]}"""

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2000,
            temperature=0.3
        )

        raw = response.choices[0].message.content.strip()

        # Clean JSON
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]

        questions = json.loads(raw)
        store["questions"] = questions
        store["answers"] = {}

        return jsonify({"success": True, "questions": questions})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/save-answer", methods=["POST"])
def save_answer():
    data = request.json
    idx = str(data.get("index"))
    answer = data.get("answer", "")
    store["answers"][idx] = answer
    return jsonify({"success": True})


@app.route("/generate-document", methods=["GET", "POST"])
def generate_document():
    if not store["questions"]:
        return jsonify({"error": "No questions generated yet"}), 400

    # Create Word document
    doc = WordDoc()

    # Title
    title = doc.add_heading("Knowledge Transfer Document", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata table
    doc.add_paragraph("")
    meta_table = doc.add_table(rows=4, cols=2)
    meta_table.style = "Table Grid"
    meta_data = [
        ("Person Leaving", store["person_name"] or "Not specified"),
        ("Role", store["role"] or "Not specified"),
        ("Project / Context", store["project"] or "Not specified"),
        ("Date", datetime.now().strftime("%d %B %Y")),
    ]
    for i, (label, value) in enumerate(meta_data):
        meta_table.rows[i].cells[0].text = label
        meta_table.rows[i].cells[1].text = value
        meta_table.rows[i].cells[0].paragraphs[0].runs[0].bold = True

    doc.add_paragraph("")

    # Documents uploaded
    doc.add_heading("Documents Reviewed", level=1)
    for d in store["documents"]:
        doc.add_paragraph(f"• {d['filename']}", style="List Bullet")

    doc.add_paragraph("")

    # Group by category
    categories = {}
    for i, q in enumerate(store["questions"]):
        cat = q["category"]
        if cat not in categories:
            categories[cat] = []
        categories[cat].append((i, q["question"]))

    # Q&A by category
    for cat, qs in categories.items():
        doc.add_heading(cat, level=1)
        for idx, question in qs:
            # Question
            q_para = doc.add_paragraph()
            q_run = q_para.add_run(f"Q: {question}")
            q_run.bold = True
            q_run.font.size = Pt(11)

            # Answer
            answer = store["answers"].get(str(idx), "").strip()
            a_para = doc.add_paragraph()
            a_run = a_para.add_run(f"A: {answer if answer else '[No answer provided]'}")
            a_run.font.size = Pt(11)
            if not answer:
                a_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

            doc.add_paragraph("")

    # Footer note
    doc.add_paragraph("")
    footer = doc.add_paragraph(
        f"This knowledge transfer document was generated on {datetime.now().strftime('%d %B %Y')} "
        f"using the KT Capture Tool."
    )
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.runs[0].font.size = Pt(9)
    footer.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Save to buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    filename = f"KT_{store['person_name'].replace(' ', '_') or 'Document'}_{datetime.now().strftime('%Y%m%d')}.docx"

    return send_file(
        buffer,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@app.route("/reset", methods=["POST"])
def reset():
    store["documents"] = []
    store["questions"] = []
    store["answers"] = {}
    store["person_name"] = ""
    store["role"] = ""
    store["project"] = ""
    return jsonify({"success": True})


if __name__ == "__main__":
    print("\n" + "="*55)
    print("  Knowledge Transfer Capture Tool")
    print("  Built with Groq + Python-docx + Flask")
    print("="*55)

    if not os.environ.get("GROQ_API_KEY"):
        print("\n  WARNING: GROQ_API_KEY not set")
        print("  Set it: set GROQ_API_KEY=gsk_...")
    else:
        print("\n  API key detected - ready to go!")

    print("\n  Open browser at: http://localhost:5000")
    print("  Press Ctrl+C to stop\n")
    print("="*55 + "\n")

    app.run(debug=False, host="0.0.0.0", port=5000)
